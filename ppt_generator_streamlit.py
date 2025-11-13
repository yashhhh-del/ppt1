import streamlit as st
import requests
import base64
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import time
import json

# Page configuration
st.set_page_config(
    page_title="AI PowerPoint Generator",
    page_icon="📊",
    layout="wide"
)

# Custom CSS
st.markdown("""
<style>
    .main-header {
        font-size: 2.5rem;
        font-weight: bold;
        color: #1f77b4;
        text-align: center;
        margin-bottom: 2rem;
    }
    .stButton>button {
        width: 100%;
        background-color: #1f77b4;
        color: white;
        font-weight: bold;
        padding: 0.5rem;
        border-radius: 5px;
    }
</style>
""", unsafe_allow_html=True)

# Header
st.markdown('<div class="main-header">📊 AI PowerPoint Generator</div>', unsafe_allow_html=True)
st.markdown("---")

# Sidebar
with st.sidebar:
    st.header("⚙️ API Configuration")
    
    claude_api_key = st.text_input("OpenRouter API Key *", type="password", 
                                    help="Required: For generating presentation content")
    
    model_choice = st.selectbox(
        "AI Model",
        [
            "Free Model (Google Gemini Flash)",
            "Free Model (Meta Llama 3.2)",
            "Free Model (Mistral 7B)",
            "Claude 3.5 Sonnet (Paid)"
        ],
        help="Try different free models if one is rate-limited"
    )
    
    if "Free" in model_choice:
        st.info("💡 Free models share rate limits. Switch models if limited.")
    
    st.info("💡 Using OpenRouter API")
    
    # Pexels API (Optional but recommended)
    pexels_api_key = st.text_input(
        "Pexels API Key (Optional)", 
        type="password",
        help="FREE! Get better topic-relevant images. Sign up at: https://www.pexels.com/api/"
    )
    
    if pexels_api_key:
        st.success("✅ Pexels connected - will search for topic-specific images!")
    
    stability_api_key = st.text_input("Stability AI API Key (Optional)", type="password", 
                                      help="Optional: For AI-generated images")
    
    st.markdown("---")
    st.markdown("### 📖 How to Use")
    st.markdown("""
    1. Enter OpenRouter API key
    2. **Optional:** Add Pexels key for better images
    3. Enter your presentation topic
    4. Click Generate!
    """)
    st.markdown("---")
    st.markdown("### 🔗 Get API Keys")
    st.markdown("🆓 [Pexels API (FREE)](https://www.pexels.com/api/)")
    st.markdown("[OpenRouter API](https://openrouter.ai/keys)")
    st.markdown("[Stability AI](https://platform.stability.ai)")

# ============ IMAGE FUNCTIONS ============

def generate_topic_search_terms(main_topic, slide_title, image_prompt):
    """Generate search terms prioritizing topic relevance"""
    search_terms = []
    
    # 1. AI's specific image prompt
    if image_prompt and image_prompt.strip():
        search_terms.append(image_prompt.strip())
    
    # 2. Topic + slide title combined
    if main_topic and slide_title:
        search_terms.append(f"{main_topic} {slide_title}")
    
    # 3. Just slide title
    if slide_title:
        search_terms.append(slide_title)
    
    # 4. Just main topic
    if main_topic:
        search_terms.append(main_topic)
    
    # Remove duplicates
    seen = set()
    unique = []
    for term in search_terms:
        lower = term.lower().strip()
        if lower and lower not in seen:
            seen.add(lower)
            unique.append(term)
    
    return unique

def get_unsplash_image(query, width=800, height=600):
    """Get image from Unsplash"""
    try:
        clean_query = query.strip().replace(' ', ',')
        url = f"https://source.unsplash.com/{width}x{height}/?{clean_query}"
        
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
        }
        
        response = requests.get(url, timeout=15, allow_redirects=True, headers=headers)
        
        if response.status_code == 200 and len(response.content) > 5000:
            # Validate image
            try:
                img = Image.open(io.BytesIO(response.content))
                if img.size[0] > 400 and img.size[1] > 300:
                    return response.content
            except:
                pass
        return None
    except:
        return None

def get_pexels_image(query, api_key):
    """Get image from Pexels API"""
    if not api_key:
        return None
    
    try:
        url = "https://api.pexels.com/v1/search"
        headers = {"Authorization": api_key}
        params = {
            "query": query,
            "per_page": 3,
            "orientation": "landscape"
        }
        
        response = requests.get(url, headers=headers, params=params, timeout=10)
        
        if response.status_code == 200:
            data = response.json()
            if data.get("photos"):
                photo = data["photos"][0]
                img_url = photo["src"]["large"]
                
                img_response = requests.get(img_url, timeout=10)
                if img_response.status_code == 200:
                    return img_response.content
        return None
    except:
        return None

def get_topic_relevant_image(main_topic, slide_title, image_prompt, pexels_key=None):
    """Get highly relevant image for the topic"""
    
    st.write(f"   🎯 Topic: {main_topic}")
    st.write(f"   📄 Slide: {slide_title}")
    
    # Generate search terms
    search_terms = generate_topic_search_terms(main_topic, slide_title, image_prompt)
    st.write(f"   🔍 Will try {len(search_terms)} search variations")
    
    # Try each search term
    for i, term in enumerate(search_terms, 1):
        st.write(f"      → Search {i}: '{term}'")
        
        # Try Pexels first if available
        if pexels_key:
            image_data = get_pexels_image(term, pexels_key)
            if image_data:
                st.write(f"      ✅ Found on Pexels!")
                return image_data
        
        # Try Unsplash
        image_data = get_unsplash_image(term)
        if image_data:
            st.write(f"      ✅ Found on Unsplash!")
            return image_data
    
    # Fallback to generic topic
    st.write(f"   🆘 Trying generic fallback...")
    fallback = main_topic.split()[0] if main_topic else "business"
    image_data = get_unsplash_image(fallback)
    if image_data:
        st.write(f"      ✅ Got fallback image")
        return image_data
    
    return None

def generate_image_stability(api_key, prompt):
    """Generate AI image using Stability AI"""
    try:
        url = "https://api.stability.ai/v2beta/stable-image/generate/core"
        
        response = requests.post(
            url,
            headers={
                "authorization": f"Bearer {api_key.strip()}",
                "accept": "image/*"
            },
            files={"none": ''},
            data={
                "prompt": f"{prompt}, professional, clean, abstract",
                "output_format": "png",
            },
        )
        
        if response.status_code == 200:
            return response.content
        return None
    except:
        return None

# ============ CONTENT GENERATION ============

def generate_content_with_claude(api_key, topic, category, slide_count, tone, audience, key_points, model_choice):
    """Generate presentation content using AI"""
    try:
        # Model selection logic
        if "Gemini" in model_choice:
            model = "google/gemini-2.0-flash-exp:free"
        elif "Llama" in model_choice:
            model = "meta-llama/llama-3.2-3b-instruct:free"
        elif "Mistral" in model_choice:
            model = "mistralai/mistral-7b-instruct:free"
        else:
            model = "anthropic/claude-3.5-sonnet"
        
        calculated_tokens = min(slide_count * 150 + 200, 1500)
        
        prompt = f"""You are an expert presentation creator. Generate a PowerPoint structure about: {topic}

Category: {category}
Slides: {slide_count}
Tone: {tone}
Audience: {audience}
Key Points: {key_points if key_points else "None"}

Return ONLY valid JSON in this format:
{{
  "slides": [
    {{
      "title": "Presentation Title",
      "bullets": [],
      "image_prompt": "{topic} title image"
    }},
    {{
      "title": "Slide Title",
      "bullets": ["Point 1", "Point 2", "Point 3"],
      "image_prompt": "specific image description related to {topic}"
    }}
  ]
}}

CRITICAL: image_prompt must be specific to {topic}. Examples:
- For "Space Exploration": use "astronaut spacewalk", "mars rover", "space station"
- For "Cooking": use "chef cooking", "fresh ingredients", "plated dish"
- For "AI": use "artificial intelligence", "neural network", "robot technology"

Make image prompts HIGHLY SPECIFIC to the topic "{topic}".
Total slides: exactly {slide_count}
Return ONLY JSON, no markdown."""

        response = requests.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {api_key.strip()}",
                "Content-Type": "application/json",
            },
            json={
                "model": model,
                "max_tokens": calculated_tokens,
                "messages": [{"role": "user", "content": prompt}]
            },
            timeout=30
        )
        
        if response.status_code == 200:
            data = response.json()
            content_text = data["choices"][0]["message"]["content"]
            
            # Clean JSON
            content_text = content_text.strip()
            if content_text.startswith("```json"):
                content_text = content_text[7:]
            if content_text.startswith("```"):
                content_text = content_text[3:]
            if content_text.endswith("```"):
                content_text = content_text[:-3]
            content_text = content_text.strip()
            
            slides_data = json.loads(content_text)
            return slides_data["slides"]
        else:
            # Enhanced error handling
            if response.status_code == 429:
                error_data = response.json()
                error_msg = error_data.get("error", {}).get("metadata", {}).get("raw", "Rate limited")
                st.error(f"⏱️ Rate Limit: Model is temporarily unavailable")
                st.info("💡 **Solutions:**\n- Wait 30-60 seconds and try again\n- Switch to a different free model above\n- Use Claude 3.5 Sonnet (paid but reliable)")
                raise Exception("Rate limit - retry needed")
            elif response.status_code == 402:
                st.error("💳 Insufficient credits! Reduce slides or add credits.")
            else:
                st.error(f"API Error: {response.text}")
            return None
            
    except json.JSONDecodeError as e:
        st.error(f"JSON parsing error: {str(e)}")
        return None
    except Exception as e:
        if "Rate limit" in str(e):
            raise  # Re-raise for retry logic
        st.error(f"Error: {str(e)}")
        return None

def generate_content_with_retry(api_key, topic, category, slide_count, tone, audience, key_points, model_choice, max_retries=3):
    """Generate content with automatic retry on rate limit"""
    for attempt in range(max_retries):
        try:
            result = generate_content_with_claude(api_key, topic, category, slide_count, tone, audience, key_points, model_choice)
            if result:
                return result
        except Exception as e:
            if "Rate limit" in str(e) or "429" in str(e):
                if attempt < max_retries - 1:
                    wait_time = (attempt + 1) * 5  # 5, 10, 15 seconds
                    st.warning(f"⏳ Rate limit hit. Retrying in {wait_time} seconds... (Attempt {attempt + 2}/{max_retries})")
                    time.sleep(wait_time)
                else:
                    st.error("❌ Rate limit persists after retries. Please:\n1. Wait 1-2 minutes\n2. Switch to different free model\n3. Use Claude model (paid)")
                    return None
            else:
                return None
    return None

# ============ POWERPOINT CREATION ============

def create_powerpoint(slides_content, theme, image_mode, stability_key, pexels_key, category, audience, topic):
    """Create PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    themes = {
        "Corporate Blue": {"bg": RGBColor(240, 248, 255), "accent": RGBColor(31, 119, 180), "text": RGBColor(0, 0, 0)},
        "Gradient Modern": {"bg": RGBColor(240, 242, 246), "accent": RGBColor(138, 43, 226), "text": RGBColor(0, 0, 0)},
        "Minimal Dark": {"bg": RGBColor(30, 30, 30), "accent": RGBColor(255, 215, 0), "text": RGBColor(255, 255, 255)},
        "Pastel Soft": {"bg": RGBColor(255, 250, 240), "accent": RGBColor(255, 182, 193), "text": RGBColor(60, 60, 60)}
    }
    
    color_scheme = themes.get(theme, themes["Corporate Blue"])
    
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    for idx, slide_data in enumerate(slides_content):
        status_text.text(f"Creating slide {idx + 1}/{len(slides_content)}...")
        progress_bar.progress((idx + 1) / len(slides_content))
        
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color_scheme["bg"]
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = slide_data["title"]
        title_frame.paragraphs[0].font.size = Pt(36 if idx == 0 else 28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = color_scheme["accent"]
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if idx == 0 else PP_ALIGN.LEFT
        
        # Bullets
        if idx > 0 and slide_data.get("bullets"):
            bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(5.5), Inches(4.5))
            text_frame = bullet_box.text_frame
            text_frame.word_wrap = True
            
            for bullet in slide_data["bullets"]:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(18)
                p.font.color.rgb = color_scheme["text"]
                p.space_after = Pt(12)
        
        # Subtitle on first slide
        if idx == 0:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = f"{category} Presentation | {audience}"
            subtitle_frame.paragraphs[0].font.size = Pt(20)
            subtitle_frame.paragraphs[0].font.color.rgb = color_scheme["text"]
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add images to content slides
        if idx > 0 and image_mode != "None":
            with st.expander(f"🖼️ Slide {idx + 1}: {slide_data['title']}", expanded=False):
                image_prompt = slide_data.get("image_prompt", "")
                image_data = None
                
                if image_mode == "AI Generated (Paid)" and stability_key:
                    st.write("   🤖 Generating AI image...")
                    image_data = generate_image_stability(stability_key, image_prompt)
                
                if not image_data and image_mode != "None":
                    st.write("   🔍 Searching for topic-relevant image...")
                    image_data = get_topic_relevant_image(
                        main_topic=topic,
                        slide_title=slide_data["title"],
                        image_prompt=image_prompt,
                        pexels_key=pexels_key
                    )
                
                if image_data:
                    try:
                        image_stream = io.BytesIO(image_data)
                        left = Inches(6.5)
                        top = Inches(2)
                        width = Inches(3)
                        pic = slide.shapes.add_picture(image_stream, left, top, width=width)
                        st.success(f"   ✅ Image added successfully!")
                    except Exception as e:
                        st.error(f"   ❌ Failed to add image: {str(e)}")
                else:
                    st.warning(f"   ⚠️ No image found")
            
            time.sleep(0.5)
    
    progress_bar.progress(1.0)
    status_text.text("✅ Presentation created!")
    
    return prs

# ============ MAIN UI ============

col1, col2 = st.columns([1, 1])

with col1:
    st.subheader("📝 Your Topic")
    topic = st.text_input("Enter Topic *", placeholder="e.g., Space Exploration, Digital Marketing, Climate Change...")
    st.caption("💡 Be specific! The more detailed your topic, the better the images will match.")
    
    category = st.selectbox("Category *", ["Business", "Pitch", "Marketing", "Technical", "Academic", "Training"])
    slide_count = st.number_input("Number of Slides *", min_value=3, max_value=15, value=6)
    tone = st.selectbox("Tone *", ["Formal", "Neutral", "Inspirational"])

with col2:
    st.subheader("🎨 Style & Images")
    audience = st.selectbox("Audience *", ["Investors", "Students", "Corporate", "Clients", "Managers"])
    theme = st.selectbox("Theme *", ["Corporate Blue", "Gradient Modern", "Minimal Dark", "Pastel Soft"])
    
    image_mode = st.selectbox(
        "Image Mode *",
        ["Free Images (Auto)", "AI Generated (Paid)", "None"],
        help="Auto: Uses Unsplash/Pexels for topic-relevant free images"
    )

st.subheader("➕ Additional Points (Optional)")
key_points = st.text_area("Key points to cover", placeholder="- Point 1\n- Point 2", height=80)

st.markdown("---")
st.info("🎯 **Topic-Specific Images**: AI will search for images that match YOUR topic!")
st.warning("💡 **Tip**: Add a FREE Pexels API key in sidebar for even better image matching!")

generate_button = st.button("🚀 Generate PowerPoint", use_container_width=True)

if generate_button:
    if not claude_api_key:
        st.error("⚠️ Enter OpenRouter API key")
    elif not topic:
        st.error("⚠️ Enter a topic")
    else:
        with st.spinner("🤖 Generating your presentation..."):
            # Use retry function
            slides_content = generate_content_with_retry(
                claude_api_key, topic, category, slide_count, 
                tone, audience, key_points, model_choice
            )
            
            if slides_content:
                st.success("✅ Content generated! Adding topic-relevant images...")
                
                prs = create_powerpoint(
                    slides_content, theme, image_mode, 
                    stability_api_key, pexels_api_key,
                    category, audience, topic
                )
                
                pptx_io = io.BytesIO()
                prs.save(pptx_io)
                pptx_io.seek(0)
                
                st.success("🎉 PowerPoint ready!")
                
                with st.expander("📄 Preview"):
                    for i, slide in enumerate(slides_content):
                        st.write(f"**Slide {i+1}:** {slide['title']}")
                        if slide.get('image_prompt'):
                            st.caption(f"   🖼️ Image: {slide['image_prompt']}")
                
                st.download_button(
                    label="📥 Download PowerPoint",
                    data=pptx_io,
                    file_name=f"{topic.replace(' ', '_')}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

st.markdown("---")
st.markdown("""
<div style='text-align: center; color: #666;'>
    <p>🎯 <strong>Images match YOUR topic automatically!</strong></p>
    <p>🆓 <strong>Get Pexels API</strong> (free) for best results: <a href="https://www.pexels.com/api/">pexels.com/api</a></p>
</div>
""", unsafe_allow_html=True)
